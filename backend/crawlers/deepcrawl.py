"""Generic URL crawler and document extractor.

Supports public HTTP(S) web pages, PDFs, CSV, Excel, DOCX, PPTX, JSON, XML,
and text files. The web crawler follows same-domain links only.
"""

import asyncio
import json
import logging
import os
import re
import tempfile

from pathlib import Path
from urllib.parse import urlparse

import fitz
import httpx
import pandas as pd
from docx import Document
from lxml import etree
from bs4 import BeautifulSoup
from pptx import Presentation


from crawl4ai import AsyncWebCrawler, BrowserConfig, CacheMode, CrawlerRunConfig
from crawl4ai.content_scraping_strategy import LXMLWebScrapingStrategy
from crawl4ai.deep_crawling import BFSDeepCrawlStrategy, BestFirstCrawlingStrategy
from crawl4ai.deep_crawling.filters import FilterChain, DomainFilter, URLPatternFilter

from crawl4ai.deep_crawling.scorers import KeywordRelevanceScorer


try:
    from storage import is_bucket_configured, make_object_key, upload_file
except ModuleNotFoundError:
    from backend.storage import is_bucket_configured, make_object_key, upload_file


logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
logger = logging.getLogger(__name__)


def extract_clean_links(raw_links, html_content: str = "", base_url: str = "", markdown_text: str = "") -> dict:
    clean_internal = []
    clean_external = []
    seen = set()

    if isinstance(raw_links, dict):
        internal_list = raw_links.get("internal", []) or []
        external_list = raw_links.get("external", []) or []
        for item in internal_list:
            href = item.get("href", "") if isinstance(item, dict) else str(item or "")
            text = item.get("text", "") if isinstance(item, dict) else href
            if href and href not in seen:
                seen.add(href)
                clean_internal.append({"href": href, "text": text or href, "type": "internal"})
        for item in external_list:
            href = item.get("href", "") if isinstance(item, dict) else str(item or "")
            text = item.get("text", "") if isinstance(item, dict) else href
            if href and href not in seen:
                seen.add(href)
                clean_external.append({"href": href, "text": text or href, "type": "external"})
    elif isinstance(raw_links, list):
        for item in raw_links:
            href = item.get("href", "") if isinstance(item, dict) else str(item or "")
            text = item.get("text", "") if isinstance(item, dict) else href
            if href and href not in seen:
                seen.add(href)
                clean_internal.append({"href": href, "text": text or href, "type": "internal"})

    if html_content and isinstance(html_content, str):
        try:
            soup = BeautifulSoup(html_content, "html.parser")
            base_netloc = urlparse(base_url).netloc if base_url else ""
            for a_tag in soup.find_all("a", href=True):
                href = a_tag.get("href", "").strip()
                if href and href not in seen and not href.startswith("javascript:") and not href.startswith("#"):
                    seen.add(href)
                    text = a_tag.get_text(strip=True) or href
                    if base_netloc and href.startswith("http") and base_netloc not in href:
                        clean_external.append({"href": href, "text": text, "type": "external"})
                    else:
                        clean_internal.append({"href": href, "text": text, "type": "internal"})
        except Exception:
            pass

    if markdown_text and isinstance(markdown_text, str):
        try:
            base_netloc = urlparse(base_url).netloc if base_url else ""
            for m in re.finditer(r'\[(.*?)\]\((https?://[^\s\)]+)\)', markdown_text):
                text, href = m.group(1).strip(), m.group(2).strip()
                if href and href not in seen:
                    seen.add(href)
                    if base_netloc and base_netloc not in href:
                        clean_external.append({"href": href, "text": text or href, "type": "external"})
                    else:
                        clean_internal.append({"href": href, "text": text or href, "type": "internal"})
        except Exception:
            pass

    return {
        "internal": clean_internal,
        "external": clean_external,
        "total_count": len(clean_internal) + len(clean_external),
    }




SUPPORTED_TYPES = {
    ".html": "html",
    ".htm": "html",
    ".php": "html",
    ".asp": "html",
    ".aspx": "html",
    ".pdf": "pdf",
    ".csv": "csv",
    ".xls": "excel",
    ".xlsx": "excel",
    ".doc": "docx",
    ".docx": "docx",
    ".ppt": "pptx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".json": "json",
    ".xml": "xml",
}


def build_response(success: bool, file_type: str, data=None, message: str = ""):
    return {
        "success": success,
        "file_type": file_type,
        "message": message,
        "data": data,
    }


def validate_url(url: str) -> None:
    parsed = urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise ValueError("Enter a valid HTTP or HTTPS URL.")


def storage_key(file_type: str, url: str, filename: str) -> str:
    return make_object_key(f"deep-crawl/{file_type}", url, filename)


def maybe_upload(local_path: str, file_type: str, url: str, filename: str):
    if not is_bucket_configured():
        return None
    return upload_file(local_path, storage_key(file_type, url, filename))


def maybe_upload_json(data: dict, file_type: str, url: str):
    if not is_bucket_configured():
        return None

    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".json", mode="w", encoding="utf-8") as handle:
            json.dump(data, handle, ensure_ascii=False, indent=2, default=str)
            temp_path = handle.name
        return upload_file(temp_path, storage_key(file_type, url, "result.json"))
    finally:
        if temp_path and os.path.exists(temp_path):
            os.remove(temp_path)


def file_type_from_content_type(content_type: str) -> str | None:
    content_type = (content_type or "").lower()

    if "text/html" in content_type:
        return "html"
    if "application/pdf" in content_type:
        return "pdf"
    if "spreadsheet" in content_type or "excel" in content_type:
        return "excel"
    if "text/csv" in content_type:
        return "csv"
    if "wordprocessingml" in content_type or "msword" in content_type:
        return "docx"
    if "presentationml" in content_type or "powerpoint" in content_type:
        return "pptx"
    if "json" in content_type:
        return "json"
    if "xml" in content_type:
        return "xml"
    if "text/plain" in content_type:
        return "txt"

    return None


async def detect_file_type(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0 (compatible; MyTrackCrawler/1.0)"}

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=20.0, headers=headers) as client:
            response = await client.head(url)
            detected = file_type_from_content_type(response.headers.get("content-type", ""))
            if detected:
                return detected
    except Exception as error:
        logger.info("HEAD request failed for %s: %s", url, error)

    return SUPPORTED_TYPES.get(Path(urlparse(url).path).suffix.lower(), "html")


async def download_file(url: str, suffix: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0 (compatible; MyTrackCrawler/1.0)"}

    async with httpx.AsyncClient(follow_redirects=True, timeout=90.0, headers=headers) as client:
        response = await client.get(url)
        response.raise_for_status()

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as handle:
            handle.write(response.content)
            return handle.name


async def extract_json(url: str):
    path = await download_file(url, ".json")
    try:
        with open(path, encoding="utf-8") as handle:
            content = json.load(handle)

        return build_response(
            True,
            "json",
            {
                "s3": maybe_upload(path, "json", url, "source.json"),
                "content": content,
            },
        )
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_xml(url: str):
    path = await download_file(url, ".xml")
    try:
        root = etree.parse(path).getroot()
        content = etree.tostring(root, pretty_print=True, encoding="unicode")

        return build_response(
            True,
            "xml",
            {
                "s3": maybe_upload(path, "xml", url, "source.xml"),
                "content": content,
            },
        )
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_pdf(url: str):
    path = await download_file(url, ".pdf")
    try:
        with fitz.open(path) as document:
            pages = [{"page": index + 1, "text": page.get_text()} for index, page in enumerate(document)]

        return build_response(
            True,
            "pdf",
            {
                "s3": maybe_upload(path, "pdf", url, "source.pdf"),
                "page_count": len(pages),
                "pages": pages,
                "text": "\n".join(page["text"] for page in pages),
            },
        )
    except Exception as error:
        logger.exception("PDF extraction failed")
        return build_response(False, "pdf", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_excel(url: str):
    path = await download_file(url, ".xlsx")
    try:
        workbook = pd.ExcelFile(path)
        sheets = {}

        for name in workbook.sheet_names:
            frame = pd.read_excel(path, sheet_name=name)
            sheets[name] = {
                "rows": len(frame),
                "columns": list(frame.columns),
                "preview": frame.head(10).fillna("").to_dict(orient="records"),
            }

        return build_response(
            True,
            "excel",
            {
                "s3": maybe_upload(path, "excel", url, "source.xlsx"),
                "sheets": sheets,
            },
        )
    except Exception as error:
        logger.exception("Excel extraction failed")
        return build_response(False, "excel", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_csv(url: str):
    path = await download_file(url, ".csv")
    try:
        frame = pd.read_csv(path)
        return build_response(
            True,
            "csv",
            {
                "s3": maybe_upload(path, "csv", url, "source.csv"),
                "rows": len(frame),
                "columns": list(frame.columns),
                "preview": frame.head(10).fillna("").to_dict(orient="records"),
            },
        )
    except Exception as error:
        logger.exception("CSV extraction failed")
        return build_response(False, "csv", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_docx(url: str):
    path = await download_file(url, ".docx")
    try:
        document = Document(path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)

        return build_response(
            True,
            "docx",
            {
                "s3": maybe_upload(path, "docx", url, "source.docx"),
                "text": text,
            },
        )
    except Exception as error:
        logger.exception("DOCX extraction failed")
        return build_response(False, "docx", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_pptx(url: str):
    path = await download_file(url, ".pptx")
    try:
        presentation = Presentation(path)
        slides = []

        for number, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            slides.append({"slide": number, "text": text})

        return build_response(
            True,
            "pptx",
            {
                "s3": maybe_upload(path, "pptx", url, "source.pptx"),
                "slides": slides,
            },
        )
    except Exception as error:
        logger.exception("PPTX extraction failed")
        return build_response(False, "pptx", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


async def extract_txt(url: str):
    path = await download_file(url, ".txt")
    try:
        with open(path, encoding="utf-8", errors="ignore") as handle:
            text = handle.read()

        return build_response(
            True,
            "txt",
            {
                "s3": maybe_upload(path, "txt", url, "source.txt"),
                "text": text,
            },
        )
    except Exception as error:
        logger.exception("TXT extraction failed")
        return build_response(False, "txt", None, str(error))
    finally:
        if os.path.exists(path):
            os.remove(path)


def markdown_text(markdown) -> str:
    return (
        getattr(markdown, "fit_markdown", None)
        or getattr(markdown, "raw_markdown", None)
        or str(markdown or "")
    )


async def extract_webpage(
    url: str,
    categories: list | None = None,
    max_depth: int = 1,
    max_pages: int = 5,
):
    """Crawl a public website and follow same-domain links using Best-First or BFS strategy."""
    browser_config = BrowserConfig(
        headless=True,
        verbose=False,
        enable_stealth=True,
        ignore_https_errors=True,
        headers={"Accept-Language": "en-US,en;q=0.9"},
        extra_args=[
            "--no-sandbox",
            "--disable-setuid-sandbox",
            "--disable-gpu",
            "--disable-dev-shm-usage",
            "--js-flags=--max-old-space-size=128",
            "--lang=en-US"
        ]
    )



    domain = urlparse(url).netloc
    base_domain = domain.replace("www.", "")
    allowed_domains = list(set([domain, base_domain, f"www.{base_domain}"]))
    url_filter = FilterChain([DomainFilter(allowed_domains=allowed_domains)])


    cleaned_categories = [c.strip() for c in (categories or []) if c and str(c).strip()]

    filters = [DomainFilter(allowed_domains=allowed_domains)]
    if cleaned_categories:
        patterns = [f"*{c.lower()}*" for c in cleaned_categories]
        filters.append(URLPatternFilter(patterns=patterns))

    url_filter = FilterChain(filters)

    logger.info(f"Extracting webpage with BFS strategy (categories: {cleaned_categories or 'all'})")
    crawl_strategy = BFSDeepCrawlStrategy(
        max_depth=min(max_depth, 2),
        max_pages=min(max_pages, 10),
        include_external=False,
        filter_chain=url_filter,
    )


    crawl_config = CrawlerRunConfig(
        deep_crawl_strategy=crawl_strategy,
        cache_mode=CacheMode.BYPASS,
        page_timeout=15_000,
        wait_until="domcontentloaded",
        verbose=False,
        semaphore_count=2,
    )


    try:
        async with asyncio.timeout(90):
            async with AsyncWebCrawler(config=browser_config) as crawler:
                results = await crawler.arun(url=url, config=crawl_config)

        if not isinstance(results, list):
            results = [results]


        pages = []
        all_links_list = []
        seen_all_links = set()

        for result in results:
            raw_links = getattr(result, "links", {})
            html_raw = getattr(result, "html", "") or ""
            md_raw = markdown_text(getattr(result, "markdown", ""))
            page_url = getattr(result, "url", url)
            parsed_links = extract_clean_links(raw_links, html_content=html_raw, base_url=page_url, markdown_text=md_raw)



            for l in parsed_links["internal"] + parsed_links["external"]:
                if l["href"] not in seen_all_links:
                    seen_all_links.add(l["href"])
                    all_links_list.append(l)

            pages.append(
                {
                    "url": getattr(result, "url", url),
                    "title": getattr(result, "title", ""),
                    "success": getattr(result, "success", False),
                    "status_code": getattr(result, "status_code", None),
                    "markdown": markdown_text(getattr(result, "markdown", "")),
                    "html": getattr(result, "html", ""),
                    "metadata": getattr(result, "metadata", {}),
                    "links": parsed_links,
                    "media": getattr(result, "media", {}),
                    "error_message": getattr(result, "error_message", None),
                }
            )


        if not pages:
            return build_response(False, "html", None, "The website returned no crawlable pages.")

        output = {
            "categories": cleaned_categories,
            "total_pages": len(pages),
            "successful_pages": sum(1 for page in pages if page["success"]),
            "all_links": all_links_list,
            "total_links_found": len(all_links_list),
            "pages": pages,
        }

        resp = build_response(
            bool(output["successful_pages"]),
            "html",
            {
                "s3": maybe_upload_json(output, "html", url),
                **output,
            },
            "" if output["successful_pages"] else "The website blocked or failed the crawl.",
        )
        resp["pages"] = pages
        resp["all_links"] = all_links_list
        resp["total_links_found"] = len(all_links_list)
        resp["categories"] = cleaned_categories
        resp["total_pages"] = len(pages)
        return resp



    except asyncio.TimeoutError:
        return build_response(False, "html", None, "Crawl timed out after 90 seconds.")
    except Exception as error:
        logger.exception("Web crawl failed for %s", url)
        return build_response(False, "html", None, str(error))


def normalize_deepcrawl_output(raw_result):
    if not raw_result or not isinstance(raw_result, dict):
        return {
            "success": False,
            "file_type": "html",
            "message": "Empty or invalid crawl result.",
            "extracted_data": {
                "pages": [],
                "total_pages": 0,
                "successful_pages": 0,
                "s3": None,
            },
        }

    # If the crawl failed, return the failure info gracefully instead of crashing
    if not raw_result.get("success", False):
        return {
            "success": False,
            "file_type": raw_result.get("file_type", "html"),
            "message": raw_result.get("message", "Crawl failed"),
            "extracted_data": {
                "pages": [],
                "total_pages": 0,
                "successful_pages": 0,
                "s3": None,
            }
        }

    payload = raw_result.get("data")
    if not isinstance(payload, dict):
        payload = raw_result

    pages = (
        payload.get("pages")
        or payload.get("results")
        or payload.get("visited_pages")
        or []
    )

    normalized_pages = []
    for item in pages:
        if not isinstance(item, dict):
            continue
        normalized_pages.append(
            {
                "url": item.get("url"),
                "title": item.get("title") or item.get("meta_title") or "Untitled",
                "content": item.get("content") or item.get("text") or item.get("markdown") or "",
                "markdown": item.get("markdown") or "",
                "html": item.get("html") or "",
                "success": item.get("success", False),
                "status": item.get("status") or item.get("status_code") or 200,
                "metadata": item.get("metadata", {}),
                "links": item.get("links", {}),
                "media": item.get("media", {}),
                "error_message": item.get("error_message"),
            }
        )

    return {
        "success": raw_result.get("success", False),
        "file_type": raw_result.get("file_type", "html"),
        "message": raw_result.get("message", ""),
        "extracted_data": {
            "categories": payload.get("categories", []),
            "pages": normalized_pages,
            "total_pages": payload.get("total_pages", len(normalized_pages)),
            "successful_pages": payload.get(
                "successful_pages",
                sum(1 for page in normalized_pages if page.get("success"))
            ),
            "s3": payload.get("s3"),
        },
    }


async def deep_crawl(
    url: str,
    categories: list | None = None,
    max_depth: int = 1,
    max_pages: int = 5,
    **kwargs
):
    """Choose the correct extractor for a public HTTP(S) URL."""
    try:
        validate_url(url)
        file_type = await detect_file_type(url)

        if file_type == "html":
            result = await extract_webpage(
                url,
                categories=categories,
                max_depth=max_depth,
                max_pages=max_pages,
            )
            return normalize_deepcrawl_output(result)

        extractors = {
            "pdf": extract_pdf,
            "excel": extract_excel,
            "csv": extract_csv,
            "docx": extract_docx,
            "pptx": extract_pptx,
            "txt": extract_txt,
            "json": extract_json,
            "xml": extract_xml,
        }

        extractor = extractors.get(file_type)
        if extractor is None:
            return build_response(False, file_type, None, f"Unsupported file type: {file_type}")

        return await extractor(url)

    except Exception as error:
        logger.exception("Crawl failed")
        return build_response(False, "unknown", None, str(error))



if __name__ == "__main__":
    async def main():
        result = await deep_crawl("https://www.geeksforgeeks.org/")
        print(json.dumps(result, indent=2, default=str))

    asyncio.run(main())
